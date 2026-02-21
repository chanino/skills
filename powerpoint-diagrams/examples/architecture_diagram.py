#!/usr/bin/env python3
"""
Architecture Diagram — 3-layer swimlane system diagram with visual polish.
Run: pip install python-pptx lxml && python architecture_diagram.py && open architecture_diagram.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Containment + flow — components live within architectural layers
#   and communicate across layer boundaries.
# Layout algorithm: Swimlane — 3 horizontal lanes (Presentation, Business Logic, Data).
#   Components are centered within their lane; edges route vertically between lanes.
# Gestalt common region: Lane background rects (semi-transparent fills with subtle
#   borders) create strong visual grouping. Components inside the same lane are
#   perceived as related even without explicit edges between them.
# Gestalt similarity: Style-based colors (primary=blue, accent=orange, success=green,
#   neutral=gray) encode component role, not layer membership.
# Z-ordering: Background rect → lane rects → title bar → connectors → component
#   shapes. Connectors render below shapes so arrows don't obscure labels.
# Scan path: Top-to-bottom (following data flow from presentation → logic → data).
#   Title bar at top anchors the viewer; lane labels (top-left of each lane) orient
#   within each region.
# Data-ink: Lane fills are functional (common region grouping), not decorative.
#   Component shadows create elevation over the lane background.

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import os

SLIDE_W, SLIDE_H = 9144000, 5143500

# --- Data ---
TITLE = "Figure 3: System Architecture"

LANES = [
    ("PRESENTATION LAYER",   "EBF3FB", "C5D9F1"),
    ("BUSINESS LOGIC LAYER", "E8F8E8", "C5E8B0"),
    ("DATA LAYER",           "FEF9E7", "F9E3A6"),
]

# (id, text, preset, lane_idx, style)
# NOTE: Using roundRect for all shapes because qlmanage Quick Look doesn't render
# flowChartProcess or flowChartMagneticDisk fills. The shapes render correctly in
# actual PowerPoint — this is a qlmanage limitation for visual QC only.
COMPONENTS = [
    ("web",    "Web App",         "roundRect",  0, "primary"),
    ("mobile", "Mobile App",      "roundRect",  0, "primary"),
    ("admin",  "Admin Dashboard", "roundRect",  0, "success"),
    ("gw",     "API Gateway",     "hexagon",    0, "accent"),
    ("auth",   "Auth Service",    "roundRect",  1, "primary"),
    ("prod",   "Product Service", "roundRect",  1, "primary"),
    ("order",  "Order Service",   "roundRect",  1, "primary"),
    ("notif",  "Notif Service",   "roundRect",  1, "accent"),
    ("udb",    "User DB",         "roundRect",  2, "neutral"),
    ("pdb",    "Product DB",      "roundRect",  2, "neutral"),
    ("odb",    "Orders DB",       "roundRect",  2, "neutral"),
    ("cache",  "Cache (Redis)",   "roundRect",  2, "accent"),
]

# (src_id, tgt_id, color)
ARROWS = [
    ("web",    "auth",  "4472C4"),
    ("mobile", "prod",  "4472C4"),
    ("admin",  "order", "70AD47"),
    ("gw",     "notif", "ED7D31"),
    ("auth",   "udb",   "5D6D7E"),
    ("prod",   "pdb",   "5D6D7E"),
    ("order",  "odb",   "5D6D7E"),
    ("notif",  "cache", "ED7D31"),
]

STYLE_COLORS = {
    "primary": ("4472C4", "2E4FA3"),
    "accent":  ("ED7D31", "C05C13"),
    "success": ("70AD47", "507C32"),
    "neutral": ("5D6D7E", "2C3E50"),
}

# Icon assignments — Iconify prefix:name keys (requires: pip install cairosvg)
ICONS = {
    "web":    "mdi:web",
    "mobile": "mdi:cellphone",
    "admin":  "mdi:shield-account",
    "gw":     "mdi:router-network",
    "auth":   "mdi:lock",
    "prod":   "mdi:package-variant",
    "order":  "mdi:cart",
    "notif":  "mdi:bell",
    "udb":    "mdi:database",
    "pdb":    "mdi:database",
    "odb":    "mdi:database",
    "cache":  "mdi:flash",
}

# --- Layout ---
TITLE_H = 500000
TITLE_GAP = 91440          # 0.1" gap between title bar and first lane
LANE_MARGIN_X = 137160     # 0.15" horizontal lane inset
BOTTOM_MARGIN = 91440      # 0.1" margin from slide bottom
COMP_W, COMP_H = 1371600, 800100   # component shapes: 1.5" x 0.875" (increased height)
COL_GAP = 304800


def center_row(n, sw, gap):
    total = n * sw + (n - 1) * gap
    left = (SLIDE_W - total) // 2
    return [left + i * (sw + gap) for i in range(n)]


def make_legend_xml(
    sid,                      # shape ID counter start
    color_entries,            # list of (label, fill_hex)
    line_entries=None,        # list of (label, color_hex, is_dashed)
    x=None, y=None,           # position (defaults to bottom-right)
    slide_w=9144000,
    slide_h=5143500,
    legend_w=1645920,         # ~1.8"
):
    """Build a legend box with color swatches and optional line-style entries.

    Returns:
        (xml_list, next_sid) — list of XML strings to append to spTree,
        and the next available shape ID.
    """
    line_entries = line_entries or []
    n_color = len(color_entries)
    n_line = len(line_entries)
    n_total = n_color + n_line

    # Sizing constants (EMU)
    row_h = 256032          # ~0.28"
    title_h = 274320        # ~0.30"
    pad_top = 45720         # 0.05"
    pad_bottom = 68580      # 0.075"
    legend_h = pad_top + title_h + n_total * row_h + pad_bottom

    margin = 137160         # 0.15" from slide edges
    if x is None:
        x = slide_w - legend_w - margin
    if y is None:
        y = slide_h - legend_h - margin

    swatch_size = 164592    # ~0.18"
    swatch_left = 91440     # 0.1" from legend left edge
    label_left = swatch_left + swatch_size + 68580  # swatch + 0.075" gap
    label_w = legend_w - label_left - 45720

    line_sample_w = 365760  # 0.4"
    line_sample_left = swatch_left

    xmls = []

    # --- Container box ---
    container_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Legend Box"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{legend_w}" cy="{legend_h}"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 3200"/></a:avLst></a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="85000"/></a:srgbClr></a:solidFill>
    <a:ln w="12700"><a:solidFill><a:srgbClr val="B0B0B0"/></a:solidFill><a:round/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="t" lIns="91440" tIns="45720" rIns="91440" bIns="0">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="l"/>
      <a:r><a:rPr lang="en-US" sz="1000" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="2C3E50"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>Legend</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''
    xmls.append(container_xml)
    sid += 1

    # --- Color swatch entries ---
    for i, (label, fill_hex) in enumerate(color_entries):
        ey = y + pad_top + title_h + i * row_h
        swatch_y = ey + (row_h - swatch_size) // 2

        swatch_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Legend Swatch {i}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x + swatch_left}" y="{swatch_y}"/><a:ext cx="{swatch_size}" cy="{swatch_size}"/></a:xfrm>
    <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 10000"/></a:avLst></a:prstGeom>
    <a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>
    <a:ln w="6350"><a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill><a:round/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''
        xmls.append(swatch_xml)
        sid += 1

        label_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Legend Label {i}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x + label_left}" y="{ey}"/><a:ext cx="{label_w}" cy="{row_h}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr" lIns="0" tIns="0" rIns="0" bIns="0">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="l"/>
      <a:r><a:rPr lang="en-US" sz="900" b="0" dirty="0">
        <a:solidFill><a:srgbClr val="2C3E50"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{label}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''
        xmls.append(label_xml)
        sid += 1

    # --- Line style entries ---
    for j, (label, color_hex, is_dashed) in enumerate(line_entries):
        ey = y + pad_top + title_h + (n_color + j) * row_h
        line_y = ey + row_h // 2
        dash_val = "dash" if is_dashed else "solid"

        line_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Legend Line {j}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x + line_sample_left}" y="{line_y}"/><a:ext cx="{line_sample_w}" cy="1"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="{line_sample_w}" b="1"/>
      <a:pathLst>
        <a:path w="{line_sample_w}" h="1">
          <a:moveTo><a:pt x="0" y="0"/></a:moveTo>
          <a:lnTo><a:pt x="{line_sample_w}" y="0"/></a:lnTo>
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="19050">
      <a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>
      <a:prstDash val="{dash_val}"/>
      <a:round/>
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''
        xmls.append(line_xml)
        sid += 1

        line_label_left = line_sample_left + line_sample_w + 68580
        line_label_w = legend_w - line_label_left - 45720

        line_label_xml = f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Legend Line Label {j}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x + line_label_left}" y="{ey}"/><a:ext cx="{line_label_w}" cy="{row_h}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr" lIns="0" tIns="0" rIns="0" bIns="0">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="l"/>
      <a:r><a:rPr lang="en-US" sz="900" b="0" dirty="0">
        <a:solidFill><a:srgbClr val="2C3E50"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{label}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''
        xmls.append(line_label_xml)
        sid += 1

    return xmls, sid


def build_routed_connector_xml(waypoints, sid, name="Connector",
                                color="888888", width=19050,
                                dash="solid", tail="triangle",
                                head="none", radius=150000):
    """Orthogonal connector with curved elbows using custom geometry."""
    xs = [p[0] for p in waypoints]
    ys = [p[1] for p in waypoints]
    min_x, min_y = min(xs), min(ys)
    max_x, max_y = max(xs), max(ys)
    bbox_w = max(max_x - min_x, 1)
    bbox_h = max(max_y - min_y, 1)

    pts = [(x - min_x, y - min_y) for x, y in waypoints]

    def direction(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        if dx > 0: return 'R'
        if dx < 0: return 'L'
        if dy > 0: return 'D'
        return 'U'

    ARC = {
        ('R', 'D'): (16200000, 5400000),   ('R', 'U'): (5400000, -5400000),
        ('L', 'D'): (16200000, -5400000),  ('L', 'U'): (5400000, 5400000),
        ('D', 'R'): (10800000, -5400000),  ('D', 'L'): (0, 5400000),
        ('U', 'R'): (10800000, 5400000),   ('U', 'L'): (0, -5400000),
    }

    path_cmds = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for i in range(1, len(pts)):
        if i < len(pts) - 1:
            prev_dir = direction(pts[i-1], pts[i])
            next_dir = direction(pts[i], pts[i+1])
            seg_before = abs(pts[i][0]-pts[i-1][0]) + abs(pts[i][1]-pts[i-1][1])
            seg_after = abs(pts[i+1][0]-pts[i][0]) + abs(pts[i+1][1]-pts[i][1])
            r = min(radius, seg_before // 2, seg_after // 2)
            bx, by = pts[i][0], pts[i][1]
            if prev_dir == 'R': bx -= r
            elif prev_dir == 'L': bx += r
            elif prev_dir == 'D': by -= r
            elif prev_dir == 'U': by += r
            path_cmds.append(f'<a:lnTo><a:pt x="{bx}" y="{by}"/></a:lnTo>')
            st, sw = ARC[(prev_dir, next_dir)]
            path_cmds.append(f'<a:arcTo wR="{r}" hR="{r}" stAng="{st}" swAng="{sw}"/>')
        else:
            path_cmds.append(f'<a:lnTo><a:pt x="{pts[i][0]}" y="{pts[i][1]}"/></a:lnTo>')

    path_data = ''.join(path_cmds)
    head_xml = f'<a:headEnd type="{head}" w="med" len="med"/>' if head != "none" else '<a:headEnd type="none"/>'
    tail_xml = f'<a:tailEnd type="{tail}" w="med" len="med"/>' if tail != "none" else '<a:tailEnd type="none"/>'

    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{bbox_w}" cy="{bbox_h}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:ahLst/>
      <a:cxnLst/>
      <a:rect l="0" t="0" r="{bbox_w}" b="{bbox_h}"/>
      <a:pathLst>
        <a:path w="{bbox_w}" h="{bbox_h}">
          {path_data}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="{dash}"/>
      <a:round/>
      {head_xml}
      {tail_xml}
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''


def resolve_icon(icon_key, size=64, tint="FFFFFF", cache_dir="/tmp/diagram_icons"):
    """Resolve an icon key to a local PNG file path.

    Resolution: local file path → Iconify API → None on failure.
    """
    os.makedirs(cache_dir, exist_ok=True)

    # Local file path
    if "/" in icon_key or "\\" in icon_key or icon_key.endswith((".png", ".svg")):
        if icon_key.endswith(".png"):
            return icon_key if os.path.exists(icon_key) else None
        if icon_key.endswith(".svg") and os.path.exists(icon_key):
            try:
                import cairosvg
                png_path = os.path.join(cache_dir, os.path.basename(icon_key) + ".png")
                with open(icon_key, "r") as f:
                    svg_data = f.read()
                if tint:
                    svg_data = svg_data.replace("currentColor", f"#{tint}")
                cairosvg.svg2png(bytestring=svg_data.encode(), write_to=png_path,
                                 output_width=size, output_height=size)
                return png_path
            except ImportError:
                return None
        return None

    # Iconify API (prefix:name)
    if ":" in icon_key:
        import urllib.request
        tint_suffix = f"_{tint}" if tint else "_orig"
        cache_name = f"{icon_key.replace(':', '_')}_{size}{tint_suffix}.png"
        png_path = os.path.join(cache_dir, cache_name)
        if os.path.exists(png_path):
            return png_path
        try:
            import cairosvg
            prefix, name = icon_key.split(":", 1)
            url = f"https://api.iconify.design/{prefix}/{name}.svg?width={size}&height={size}"
            req = urllib.request.Request(url, headers={"User-Agent": "python-pptx-diagrams/1.0"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                svg_data = resp.read().decode("utf-8")
            if tint:
                svg_data = svg_data.replace("currentColor", f"#{tint}")
            cairosvg.svg2png(bytestring=svg_data.encode(), write_to=png_path,
                             output_width=size, output_height=size)
            return png_path
        except ImportError:
            print(f"  Warning: cairosvg not installed, skipping icon '{icon_key}'")
            return None
        except Exception as e:
            print(f"  Warning: Failed to resolve icon '{icon_key}': {e}")
            return None

    return None


def enrich_shape_with_icon(slide, shape_name, icon_path, icon_size_inches=0.35):
    """Add an icon image above the text label inside a named shape."""
    from pptx.util import Inches

    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break
    if target is None:
        return False

    icon_w = icon_h = Inches(icon_size_inches)
    icon_left = target.left + (target.width - icon_w) // 2
    icon_top = target.top + Emu(91440)  # 0.1" from top of shape

    try:
        slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_w, icon_h)
    except Exception as e:
        print(f"  Warning: Failed to add icon to '{shape_name}': {e}")
        return False

    # Push text below icon
    if target.has_text_frame:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        body_pr = target.text_frame._txBody.find(f"{{{ns}}}bodyPr")
        if body_pr is not None:
            body_pr.set("tIns", "365760")  # 0.4" top inset
            body_pr.set("anchor", "b")     # anchor text to bottom

    return True


def enrich_pptx_with_icons(pptx_path, icon_map, output_path=None,
                            icon_size=64, icon_size_inches=0.35,
                            tint="FFFFFF", slide_index=0):
    """Post-process a PPTX to add icons to named shapes."""
    prs_icons = Presentation(pptx_path)
    slide = prs_icons.slides[slide_index]
    results = []

    for shape_name, icon_key in icon_map.items():
        icon_path = resolve_icon(icon_key, size=icon_size, tint=tint)
        if icon_path:
            ok = enrich_shape_with_icon(slide, shape_name, icon_path, icon_size_inches)
            results.append((shape_name, ok))
        else:
            results.append((shape_name, False))

    prs_icons.save(output_path or pptx_path)
    return results


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # Background
    bg_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="F5F6FA"/></a:solidFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
    </p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml.strip()))
    sid += 1

    # Lane backgrounds
    usable_h = SLIDE_H - TITLE_H - TITLE_GAP - BOTTOM_MARGIN
    lane_h = usable_h // len(LANES)
    lane_w = SLIDE_W - 2 * LANE_MARGIN_X
    lane_y_positions = []

    for i, (label, fill, border) in enumerate(LANES):
        y = TITLE_H + TITLE_GAP + i * lane_h
        lane_y_positions.append(y)
        lane_xml = f'''
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr><p:cNvPr id="{sid}" name="Lane {label}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
          <p:spPr>
            <a:xfrm><a:off x="{LANE_MARGIN_X}" y="{y}"/><a:ext cx="{lane_w}" cy="{lane_h}"/></a:xfrm>
            <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
            <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
            <a:ln w="12700"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill></a:ln>
          </p:spPr>
          <p:txBody>
            <a:bodyPr wrap="square" anchor="t" lIns="91440" tIns="91440"><a:normAutofit/></a:bodyPr>
            <a:lstStyle/>
            <a:p><a:pPr algn="l"/>
              <a:r><a:rPr lang="en-US" sz="900" b="1" dirty="0">
                <a:solidFill><a:srgbClr val="888888"/></a:solidFill>
                <a:latin typeface="Calibri"/>
              </a:rPr><a:t>{label}</a:t></a:r>
            </a:p>
          </p:txBody>
        </p:sp>'''
        spTree.append(etree.fromstring(lane_xml.strip()))
        sid += 1

    # Title bar
    title_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="TitleBar"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{LANE_MARGIN_X}" y="0"/><a:ext cx="{lane_w}" cy="{TITLE_H}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="1F3864"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="16294A"/></a:gs>
          </a:gsLst>
          <a:lin ang="0" scaled="0"/>
        </a:gradFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="2200" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{TITLE}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''
    spTree.append(etree.fromstring(title_xml.strip()))
    sid += 1

    # Compute component positions
    lane_comps = {}
    for cid, text, preset, lane_idx, style in COMPONENTS:
        lane_comps.setdefault(lane_idx, []).append((cid, text, preset, style))

    positions = {}  # cid -> (x, y, cx, cy)
    for lane_idx, comps in lane_comps.items():
        xs = center_row(len(comps), COMP_W, COL_GAP)
        ly = lane_y_positions[lane_idx]
        comp_y = ly + (lane_h - COMP_H) // 2
        for i, (cid, text, preset, style) in enumerate(comps):
            positions[cid] = (xs[i], comp_y, COMP_W, COMP_H)

    # Arrows (render before shapes) — orthogonal routing
    for src_id, tgt_id, color in ARROWS:
        sx, sy, scx, scy = positions[src_id]
        tx, ty, tcx, tcy = positions[tgt_id]
        x1 = sx + scx // 2    # source bottom center
        y1 = sy + scy
        x2 = tx + tcx // 2    # target top center
        y2 = ty
        mid_y = y1 + (y2 - y1) // 2
        if x1 == x2:
            waypoints = [(x1, y1), (x2, y2)]
        else:
            waypoints = [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]
        xml = build_routed_connector_xml(
            waypoints, sid, f"Arrow {sid}",
            color=color, width=19050, tail="stealth", radius=120000)
        spTree.append(etree.fromstring(xml))
        sid += 1

    # Component shapes (on top)
    for cid, text, preset, lane_idx, style in COMPONENTS:
        x, y, cx, cy = positions[cid]
        fill, border = STYLE_COLORS[style]
        comp_xml = f'''
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="{sid}" name="{cid}"/>
            <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
            <p:nvPr/>
          </p:nvSpPr>
          <p:spPr>
            <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
            <a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>
            <a:gradFill>
              <a:gsLst>
                <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
                <a:gs pos="100000"><a:srgbClr val="{border}"/></a:gs>
              </a:gsLst>
              <a:lin ang="5400000" scaled="0"/>
            </a:gradFill>
            <a:ln w="25400"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
            <a:effectLst>
              <a:outerShdw blurRad="50800" dist="38100" dir="5400000" algn="tl" rotWithShape="0">
                <a:srgbClr val="000000"><a:alpha val="35000"/></a:srgbClr>
              </a:outerShdw>
            </a:effectLst>
          </p:spPr>
          <p:txBody>
            <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
            <a:lstStyle/>
            <a:p><a:pPr algn="ctr"/>
              <a:r><a:rPr lang="en-US" sz="1200" b="1" dirty="0">
                <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
                <a:latin typeface="Calibri"/>
              </a:rPr><a:t>{text}</a:t></a:r>
            </a:p>
          </p:txBody>
        </p:sp>'''
        spTree.append(etree.fromstring(comp_xml.strip()))
        sid += 1

    # Legend (last in z-order — renders on top)
    legend_entries = [
        ("Interface", STYLE_COLORS["primary"][0]),
        ("Gateway/Async", STYLE_COLORS["accent"][0]),
        ("Admin", STYLE_COLORS["success"][0]),
        ("Data Store", STYLE_COLORS["neutral"][0]),
    ]
    line_entries = [
        ("Data flow", "4472C4", False),
        ("Async", "ED7D31", True),
    ]
    # Bottom-right is occupied by Cache — place legend bottom-left instead
    # Narrower width (1.2") to fit left of User DB (which starts at x=1.5")
    legend_xmls, sid = make_legend_xml(sid, legend_entries, line_entries,
                                        x=137160, legend_w=1097280)  # 0.15" left, 1.2" wide
    for xml_str in legend_xmls:
        spTree.append(etree.fromstring(xml_str))

    output = "architecture_diagram.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")

    # --- Icon enrichment (optional post-processing) ---
    if ICONS:
        try:
            results = enrich_pptx_with_icons(output, ICONS, tint="FFFFFF")
            added = sum(1 for _, ok in results if ok)
            print(f"Icons: {added}/{len(ICONS)} added successfully")
        except ImportError:
            print("Note: Install cairosvg for icon support (pip install cairosvg)")


if __name__ == "__main__":
    main()
