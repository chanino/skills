#!/usr/bin/env python3
"""
Branching Flowchart — Decision-based flow with orthogonal routing and curved elbows.
Demonstrates: straight arrows, L-shaped paths, Z-shaped paths, loop-back paths,
YES/NO decision branches, and arrowheads on curved connectors.

Run: pip install python-pptx lxml && python branching_flowchart.py && open branching_flowchart.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Decision logic — a process with conditional branching and a
#   retry loop. The primary question is "what happens at each gate?"
# Layout algorithm: Layered flow with 2 rows. Row 1 = initial path (START →
#   VALIDATE → VALID?). Row 2 = YES path (PROCESS → COMPLETE? → DONE).
#   REJECT branches right from VALID? (same row, extending the scan line).
# Double coding: YES paths use green color + solid stroke. NO paths use red
#   color + dashed stroke. This ensures decision outcomes are distinguishable
#   even in grayscale or for colorblind viewers.
# Edge bundling: The COMPLETE?→PROCESS loop-back routes below the main flow
#   (down, left, up) to avoid crossing the primary left-to-right path. Dashed
#   stroke further separates it visually from forward-flow arrows.
# Gestalt continuity: All connectors use orthogonal routing with curved elbows
#   (consistent radius R=150000). Smooth, predictable paths reduce cognitive load.
# Scan path: Left-to-right for the happy path, then decision branches pull the
#   eye down (YES) or right (NO). Loop-back is visually subordinate (dashed, below).
# Data-ink: Labels (YES/NO) are placed adjacent to decision exits, not on the
#   arrows themselves, avoiding clutter while maintaining clarity.

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import os

SLIDE_W, SLIDE_H = 9144000, 5143500
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# --- Colors ---
CLR_TERM = ("1F3864", "16294A")    # terminators
CLR_PROC = ("4472C4", "2E4FA3")    # process
CLR_DEC = ("ED7D31", "C05C13")     # decision
CLR_ALT = ("70AD47", "507C32")     # alternate path
CLR_ARROW = "888888"
CLR_YES = "27AE60"
CLR_NO = "E74C3C"

# --- Layout ---
SW, SH = 1508760, 800100       # shape width/height (1.65" x 0.875")
DW, DH = 1508760, 1028700      # decision width/height (1.65" x 1.125")
GAP_H, GAP_V = 274320, 457200  # horizontal/vertical gaps (0.3" h, 0.5" v)
R = 150000                       # corner radius

# 4-column grid: REJECT shares col4 with DONE (row1 vs row2)
_TOTAL_W = 4 * SW + 3 * GAP_H
_LEFT = (SLIDE_W - _TOTAL_W) // 2
MARGIN_TOP = 700000


# ═══════════════════════════════════════════════════════════════════
# Orthogonal connector with curved elbows (custom geometry)
# ═══════════════════════════════════════════════════════════════════

def build_routed_connector_xml(waypoints, sid, name="Connector",
                                color="888888", width=19050,
                                dash="solid", tail="triangle",
                                head="none", radius=150000):
    xs = [p[0] for p in waypoints]
    ys = [p[1] for p in waypoints]
    min_x, min_y = min(xs), min(ys)
    max_x, max_y = max(xs), max(ys)
    bbox_w = max(max_x - min_x, 1)
    bbox_h = max(max_y - min_y, 1)

    pts = [(x - min_x, y - min_y) for x, y in waypoints]

    def direction(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        if dx > 0: return 'R'
        if dx < 0: return 'L'
        if dy > 0: return 'D'
        return 'U'

    ARC = {
        ('R', 'D'): (16200000, 5400000),   ('R', 'U'): (5400000, -5400000),
        ('L', 'D'): (16200000, -5400000),  ('L', 'U'): (5400000, 5400000),
        ('D', 'R'): (10800000, -5400000),  ('D', 'L'): (0, 5400000),
        ('U', 'R'): (10800000, 5400000),   ('U', 'L'): (0, -5400000),
    }

    path_cmds = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']

    for i in range(1, len(pts)):
        if i < len(pts) - 1:
            prev_dir = direction(pts[i-1], pts[i])
            next_dir = direction(pts[i], pts[i+1])
            seg_before = abs(pts[i][0]-pts[i-1][0]) + abs(pts[i][1]-pts[i-1][1])
            seg_after = abs(pts[i+1][0]-pts[i][0]) + abs(pts[i+1][1]-pts[i][1])
            r = min(radius, seg_before // 2, seg_after // 2)

            bx, by = pts[i][0], pts[i][1]
            if prev_dir == 'R': bx -= r
            elif prev_dir == 'L': bx += r
            elif prev_dir == 'D': by -= r
            elif prev_dir == 'U': by += r
            path_cmds.append(f'<a:lnTo><a:pt x="{bx}" y="{by}"/></a:lnTo>')

            st, sw = ARC[(prev_dir, next_dir)]
            path_cmds.append(f'<a:arcTo wR="{r}" hR="{r}" stAng="{st}" swAng="{sw}"/>')
        else:
            path_cmds.append(
                f'<a:lnTo><a:pt x="{pts[i][0]}" y="{pts[i][1]}"/></a:lnTo>'
            )

    path_data = ''.join(path_cmds)
    head_xml = f'<a:headEnd type="{head}" w="med" len="med"/>' if head != "none" else '<a:headEnd type="none"/>'
    tail_xml = f'<a:tailEnd type="{tail}" w="med" len="med"/>' if tail != "none" else '<a:tailEnd type="none"/>'

    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{bbox_w}" cy="{bbox_h}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:ahLst/>
      <a:cxnLst/>
      <a:rect l="0" t="0" r="{bbox_w}" b="{bbox_h}"/>
      <a:pathLst>
        <a:path w="{bbox_w}" h="{bbox_h}">
          {path_data}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="{dash}"/>
      <a:round/>
      {head_xml}
      {tail_xml}
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''


# ═══════════════════════════════════════════════════════════════════
# Shape + label helpers
# ═══════════════════════════════════════════════════════════════════

def make_shape_xml(sid, name, text, preset, x, y, cx, cy, fill, border):
    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{border}"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:ln w="19050"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
    <a:effectLst>
      <a:outerShdw blurRad="50800" dist="38100" dir="5400000" algn="tl" rotWithShape="0">
        <a:srgbClr val="000000"><a:alpha val="35000"/></a:srgbClr>
      </a:outerShdw>
    </a:effectLst>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="en-US" sz="1200" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{text}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def make_label_xml(sid, text, x, y, color="666666", sz=900):
    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr><p:cNvPr id="{sid}" name="Label {sid}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="342900" cy="304800"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="en-US" sz="{sz}" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{text}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # Background
    bg_xml = f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="F8FAFF"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="EEF2FA"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml))
    sid += 1

    # Title
    title_xml = f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr><p:cNvPr id="{sid}" name="Title"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="152400"/><a:ext cx="8229600" cy="457200"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="en-US" sz="2200" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="1F3864"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>Figure 4: Request Processing Flow</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(title_xml))
    sid += 1

    # ── Shape positions ──────────────────────────────────────────
    # 4-column grid layout:
    # Row 1: START(col1) → VALIDATE(col2) → VALID?(col3) → REJECT(col4)
    # Row 2:               PROCESS(col2) → COMPLETE?(col3) → DONE(col4)
    # Loop-back from COMPLETE?=NO to PROCESS

    col1_x = _LEFT
    col2_x = _LEFT + (SW + GAP_H)
    col3_x = _LEFT + 2 * (SW + GAP_H)
    col4_x = _LEFT + 3 * (SW + GAP_H)

    row1_y = MARGIN_TOP
    row2_y = row1_y + DH + GAP_V  # use DH since diamonds are tallest in row 1

    shapes = {
        # (text, preset, x, y, cx, cy, fill, border)
        "start":    ("START",     "roundRect",          col1_x, row1_y + (DH - SH) // 2, SW, SH, *CLR_TERM),
        "validate": ("Validate\nInput", "roundRect",    col2_x, row1_y + (DH - SH) // 2, SW, SH, *CLR_PROC),
        "valid":    ("Valid?",    "flowChartDecision",   col3_x, row1_y, SW, DH, *CLR_DEC),
        "process":  ("Process\nRequest", "roundRect",   col2_x, row2_y + (DH - SH) // 2, SW, SH, *CLR_PROC),
        "complete": ("Complete?", "flowChartDecision",   col3_x, row2_y, SW, DH, *CLR_DEC),
        "done":     ("DONE",      "roundRect",          col4_x, row2_y + (DH - SH) // 2, SW, SH, *CLR_TERM),
        "reject":   ("Reject\nRequest", "roundRect",    col4_x, row1_y + (DH - SH) // 2, SW, SH, *CLR_ALT),
    }

    # Helper to get center/edge points
    def cx_of(key): return shapes[key][2] + shapes[key][4] // 2
    def cy_of(key): return shapes[key][3] + shapes[key][5] // 2
    def right_of(key): return shapes[key][2] + shapes[key][4]
    def left_of(key): return shapes[key][2]
    def top_of(key): return shapes[key][3]
    def bot_of(key): return shapes[key][3] + shapes[key][5]

    # ── Connectors (render before shapes for z-order) ──

    # 1. START → VALIDATE (straight horizontal)
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(right_of("start"), cy_of("start")),
         (left_of("validate"), cy_of("validate"))],
        sid, "Arrow Start-Validate", CLR_ARROW, tail="triangle", radius=R)))
    sid += 1

    # 2. VALIDATE → VALID? (straight horizontal)
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(right_of("validate"), cy_of("validate")),
         (left_of("valid"), cy_of("valid"))],
        sid, "Arrow Validate-Valid", CLR_ARROW, tail="triangle", radius=R)))
    sid += 1

    # 3. VALID? → REJECT (YES=right exit → right, NO path)
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(right_of("valid"), cy_of("valid")),
         (left_of("reject"), cy_of("reject"))],
        sid, "Arrow Valid-Reject", CLR_NO, tail="triangle", radius=R)))
    sid += 1

    # 4. VALID? → PROCESS (YES=bottom exit → down then left, L-shape)
    mid_y = bot_of("valid") + (top_of("process") - bot_of("valid")) // 2
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(cx_of("valid"), bot_of("valid")),
         (cx_of("valid"), mid_y),
         (cx_of("process"), mid_y),
         (cx_of("process"), top_of("process"))],
        sid, "Arrow Valid-Process", CLR_YES, tail="triangle", radius=R)))
    sid += 1

    # 5. PROCESS → COMPLETE? (straight horizontal)
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(right_of("process"), cy_of("process")),
         (left_of("complete"), cy_of("complete"))],
        sid, "Arrow Process-Complete", CLR_ARROW, tail="triangle", radius=R)))
    sid += 1

    # 6. COMPLETE? → DONE (YES, straight horizontal)
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(right_of("complete"), cy_of("complete")),
         (left_of("done"), cy_of("done"))],
        sid, "Arrow Complete-Done", CLR_YES, tail="triangle", radius=R)))
    sid += 1

    # 7. COMPLETE? → PROCESS (NO, loop-back: down, left, up)
    loop_y = bot_of("complete") + 300000
    spTree.append(etree.fromstring(build_routed_connector_xml(
        [(cx_of("complete"), bot_of("complete")),
         (cx_of("complete"), loop_y),
         (cx_of("process"), loop_y),
         (cx_of("process"), bot_of("process"))],
        sid, "Arrow Complete-Process Loop", CLR_NO, dash="dash",
        tail="triangle", radius=R)))
    sid += 1

    # ── Shapes (on top of arrows) ──
    for key, (text, preset, x, y, cx, cy, fill, border) in shapes.items():
        xml = make_shape_xml(sid, key, text, preset, x, y, cx, cy, fill, border)
        spTree.append(etree.fromstring(xml))
        sid += 1

    # ── YES/NO labels ──
    # "NO" label on Valid? → Reject (centered in gap above arrow)
    gap_cx = right_of("valid") + (left_of("reject") - right_of("valid")) // 2
    spTree.append(etree.fromstring(make_label_xml(
        sid, "NO", gap_cx - 171450, cy_of("valid") - 300000, CLR_NO, 1200)))
    sid += 1

    # "YES" label on Valid? → Process (left of the down arrow, in row gap)
    spTree.append(etree.fromstring(make_label_xml(
        sid, "YES", cx_of("valid") + 100000, bot_of("valid") + 20000, CLR_YES, 1200)))
    sid += 1

    # "YES" label on Complete? → Done (centered in gap above arrow)
    gap_cx2 = right_of("complete") + (left_of("done") - right_of("complete")) // 2
    spTree.append(etree.fromstring(make_label_xml(
        sid, "YES", gap_cx2 - 171450, cy_of("complete") - 300000, CLR_YES, 1200)))
    sid += 1

    # "NO" label on Complete? loop-back (right of bottom exit)
    spTree.append(etree.fromstring(make_label_xml(
        sid, "NO", cx_of("complete") + 100000, bot_of("complete") + 20000, CLR_NO, 1200)))
    sid += 1

    output = "branching_flowchart.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")


if __name__ == "__main__":
    main()
