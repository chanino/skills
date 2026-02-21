#!/usr/bin/env python3
"""
Hub-and-Spoke Diagram — API Gateway at center with 6 services radiating outward.
Run: pip install python-pptx lxml && python hub_spoke_diagram.py && open hub_spoke_diagram.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Radial/hub-and-spoke — a central coordinator (API Gateway)
#   with dependent services arranged symmetrically around it.
# Layout algorithm: Hub-and-spoke with trigonometric positioning. Spokes are
#   placed at equal angular intervals around the hub center, producing a
#   balanced radial layout that emphasizes the hub's centrality.
# Gestalt similarity: Services share the same shape (roundRect) and size;
#   color encodes service category (blue=core, green=data, orange=infrastructure).
#   The hub is visually dominant (larger, hexagon, darker color).
# Gestalt proximity: All spokes are equidistant from the hub, reinforcing that
#   they are peers in relationship to the central coordinator.
# Scan path: Center-outward. The hub (largest, darkest) captures attention first,
#   then the eye radiates to each spoke. No left-to-right bias.
# Data-ink: Straight connectors from hub to spoke — the simplest encoding for
#   "connected to center." No bent/orthogonal routing needed since spokes are
#   arranged to avoid edge crossings by construction.
# Double coding: Hub uses a distinct shape (hexagon) in addition to a distinct
#   color, so it remains identifiable even without color.

import math
import os
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

SLIDE_W, SLIDE_H = 9144000, 5143500
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# --- Data ---
TITLE = "Figure 5: API Gateway Architecture"

HUB = ("gateway", "API\nGateway", "hexagon", "1F3864", "16294A")

SPOKES = [
    ("auth",    "Auth\nService",    "roundRect", "4472C4", "2E4FA3"),
    ("users",   "User\nService",    "roundRect", "4472C4", "2E4FA3"),
    ("orders",  "Order\nService",   "roundRect", "70AD47", "507C32"),
    ("catalog", "Catalog\nService", "roundRect", "70AD47", "507C32"),
    ("notify",  "Notification\nService", "roundRect", "ED7D31", "C05C13"),
    ("monitor", "Monitoring\nService",   "roundRect", "ED7D31", "C05C13"),
]

# --- Layout constants ---
HUB_W, HUB_H = 1828800, 914400       # 2.0" x 1.0" (1.3x primary)
SPOKE_W, SPOKE_H = 1371600, 685800   # 1.5" x 0.75" (primary size)
RADIUS = 1828800                       # 2.0" (reduced from 2.3" to fit within slide bounds)
TITLE_H = 457200
CENTER_X = SLIDE_W // 2
# Vertically center the hub+spokes in the space below title, with clearance
_TOP_CLEARANCE = TITLE_H + 228600     # 0.25" below title
_BOT_CLEARANCE = SLIDE_H - 114300     # 0.125" above bottom
CENTER_Y = (_TOP_CLEARANCE + _BOT_CLEARANCE) // 2


def make_shape_xml(sid, name, text, preset, x, y, cx, cy, fill, border,
                   font_sz=1100, shadow=True):
    """Shape with gradient fill, optional shadow, and centered text."""
    shadow_xml = ""
    if shadow:
        shadow_xml = """<a:effectLst>
          <a:outerShdw blurRad="50800" dist="38100" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="000000"><a:alpha val="35000"/></a:srgbClr>
          </a:outerShdw>
        </a:effectLst>"""

    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{border}"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:ln w="19050"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
    {shadow_xml}
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="en-US" sz="{font_sz}" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{text}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def make_straight_arrow_xml(sid, x1, y1, x2, y2, color="888888", width=19050):
    """Straight connector arrow between two points."""
    # Compute bounding box
    min_x, min_y = min(x1, x2), min(y1, y2)
    cx = abs(x2 - x1) or 1
    cy = abs(y2 - y1) or 1
    flip_h = '1' if x2 < x1 else '0'
    flip_v = '1' if y2 < y1 else '0'

    return f'''<p:cxnSp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvCxnSpPr>
    <p:cNvPr id="{sid}" name="Arrow {sid}"/>
    <p:cNvCxnSpPr/>
    <p:nvPr/>
  </p:nvCxnSpPr>
  <p:spPr>
    <a:xfrm flipH="{flip_h}" flipV="{flip_v}">
      <a:off x="{min_x}" y="{min_y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
    <a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="solid"/>
      <a:round/>
      <a:headEnd type="none"/>
      <a:tailEnd type="triangle" w="med" len="med"/>
    </a:ln>
  </p:spPr>
</p:cxnSp>'''


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # Background
    bg_xml = f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:gradFill>
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="F8FAFF"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="EEF2FA"/></a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml))
    sid += 1

    # Title
    title_xml = f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:nvSpPr><p:cNvPr id="{sid}" name="Title"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="457200" y="152400"/><a:ext cx="8229600" cy="400050"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="en-US" sz="2200" b="1" dirty="0">
        <a:solidFill><a:srgbClr val="1F3864"/></a:solidFill>
        <a:latin typeface="Calibri"/>
      </a:rPr><a:t>{TITLE}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''
    spTree.append(etree.fromstring(title_xml))
    sid += 1

    # Compute spoke positions using trigonometry
    # Start from top (270 degrees / -90 degrees) and go clockwise
    n = len(SPOKES)
    spoke_positions = {}  # key -> (center_x, center_y, left, top)
    for i, (key, text, preset, fill, border) in enumerate(SPOKES):
        angle_rad = math.radians(-90 + i * (360 / n))
        cx = int(CENTER_X + RADIUS * math.cos(angle_rad))
        cy = int(CENTER_Y + RADIUS * math.sin(angle_rad))
        left = cx - SPOKE_W // 2
        top = cy - SPOKE_H // 2
        spoke_positions[key] = (cx, cy, left, top)

    # Hub position
    hub_left = CENTER_X - HUB_W // 2
    hub_top = CENTER_Y - HUB_H // 2

    # Arrows (render before shapes for z-order)
    for key in spoke_positions:
        spoke_cx, spoke_cy, _, _ = spoke_positions[key]

        # Arrow from hub edge to spoke edge (along the radial line)
        # Compute intersection points with shape boundaries (approximated)
        dx = spoke_cx - CENTER_X
        dy = spoke_cy - CENTER_Y
        dist = math.sqrt(dx * dx + dy * dy)
        if dist == 0:
            continue
        ux, uy = dx / dist, dy / dist  # unit vector

        # Start: hub edge (approximate using ellipse for hexagon)
        hub_rx, hub_ry = HUB_W // 2, HUB_H // 2
        start_x = int(CENTER_X + hub_rx * ux * 0.85)
        start_y = int(CENTER_Y + hub_ry * uy * 0.85)

        # End: spoke edge (approximate using rectangle)
        spoke_rx, spoke_ry = SPOKE_W // 2, SPOKE_H // 2
        # Find intersection with rectangle
        if abs(ux) * spoke_ry > abs(uy) * spoke_rx:
            # Hits left/right side
            t = spoke_rx / abs(ux) if abs(ux) > 0.001 else spoke_ry / abs(uy)
        else:
            # Hits top/bottom side
            t = spoke_ry / abs(uy) if abs(uy) > 0.001 else spoke_rx / abs(ux)
        end_x = int(spoke_cx - ux * t)
        end_y = int(spoke_cy - uy * t)

        xml = make_straight_arrow_xml(sid, start_x, start_y, end_x, end_y,
                                      color="AAAAAA", width=19050)
        spTree.append(etree.fromstring(xml))
        sid += 1

    # Hub shape (on top of arrows)
    hub_key, hub_text, hub_preset, hub_fill, hub_border = HUB
    xml = make_shape_xml(sid, hub_key, hub_text, hub_preset,
                         hub_left, hub_top, HUB_W, HUB_H,
                         hub_fill, hub_border, font_sz=1400)
    spTree.append(etree.fromstring(xml))
    sid += 1

    # Spoke shapes (on top of arrows)
    for key, text, preset, fill, border in SPOKES:
        _, _, left, top = spoke_positions[key]
        xml = make_shape_xml(sid, key, text, preset,
                             left, top, SPOKE_W, SPOKE_H,
                             fill, border, font_sz=1100)
        spTree.append(etree.fromstring(xml))
        sid += 1

    output = "hub_spoke_diagram.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")


if __name__ == "__main__":
    main()
