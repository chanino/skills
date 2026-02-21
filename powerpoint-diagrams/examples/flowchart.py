#!/usr/bin/env python3
"""
Flowchart — 5-step horizontal process flow with gradient fills and shadows.
Run: pip install python-pptx lxml && python flowchart.py && open flowchart.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Sequential flow — steps executed in order.
# Layout algorithm: Horizontal flow (left-to-right scan path).
# Gestalt proximity: Equal GAP between shapes creates a uniform rhythm; the reader
#   perceives them as a single ordered sequence rather than isolated items.
# Gestalt similarity: Role-based colors — terminators (dark navy) vs. process (blue)
#   vs. decision (orange) — let the viewer distinguish shape function without reading labels.
# Scan path: Left-to-right, matching LTR reading order. Entry = START (leftmost),
#   terminal = DONE (rightmost, same dark navy signals closure).
# Data-ink: Gradients encode depth (top-lighter = 3D cue), shadows create
#   figure-ground separation. No decorative effects.

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import os

SLIDE_W, SLIDE_H = 9144000, 5143500

# --- Data ---
TITLE = "Figure 1: Project Execution Process"
STEPS = [
    ("START",   "roundRect",          "1F3864", "16294A"),
    ("PLAN",    "roundRect",          "4472C4", "2E4FA3"),
    ("REVIEW?", "flowChartDecision",  "ED7D31", "C05C13"),
    ("EXECUTE", "roundRect",          "4472C4", "2E4FA3"),
    ("DONE",    "roundRect",          "1F3864", "16294A"),
]

# --- Layout constants ---
SW, GAP = 1554480, 182880   # 1.7" shapes, 0.2" gaps
DIMS = {"roundRect": 914400, "flowChartDecision": 1143000}
N = len(STEPS)
total_w = N * SW + (N - 1) * GAP
LEFT = (SLIDE_W - total_w) // 2
BASE_CY = 1143000  # tallest shape height (decision diamond)
BASE_Y = (SLIDE_H - BASE_CY) // 2 + 114300  # shift down slightly to balance with title


def make_shape_xml(sid, name, text, preset, x, y, cx, cy, fill, border):
    """Create a shape with gradient fill, shadow, and centered text."""
    return f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{sid}" name="{name}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
        <a:prstGeom prst="{preset}"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="{border}"/></a:gs>
          </a:gsLst>
          <a:lin ang="5400000" scaled="0"/>
        </a:gradFill>
        <a:ln w="19050"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
        <a:effectLst>
          <a:outerShdw blurRad="50800" dist="38100" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="000000"><a:alpha val="35000"/></a:srgbClr>
          </a:outerShdw>
        </a:effectLst>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="1400" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{text}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''


def build_routed_connector_xml(waypoints, sid, name="Connector",
                                color="888888", width=19050,
                                dash="solid", tail="triangle",
                                head="none", radius=150000):
    """Orthogonal connector with curved elbows using custom geometry."""
    xs = [p[0] for p in waypoints]
    ys = [p[1] for p in waypoints]
    min_x, min_y = min(xs), min(ys)
    max_x, max_y = max(xs), max(ys)
    bbox_w = max(max_x - min_x, 1)
    bbox_h = max(max_y - min_y, 1)

    pts = [(x - min_x, y - min_y) for x, y in waypoints]

    def direction(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        if dx > 0: return 'R'
        if dx < 0: return 'L'
        if dy > 0: return 'D'
        return 'U'

    ARC = {
        ('R', 'D'): (16200000, 5400000),   ('R', 'U'): (5400000, -5400000),
        ('L', 'D'): (16200000, -5400000),  ('L', 'U'): (5400000, 5400000),
        ('D', 'R'): (10800000, -5400000),  ('D', 'L'): (0, 5400000),
        ('U', 'R'): (10800000, 5400000),   ('U', 'L'): (0, -5400000),
    }

    path_cmds = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for i in range(1, len(pts)):
        if i < len(pts) - 1:
            prev_dir = direction(pts[i-1], pts[i])
            next_dir = direction(pts[i], pts[i+1])
            seg_before = abs(pts[i][0]-pts[i-1][0]) + abs(pts[i][1]-pts[i-1][1])
            seg_after = abs(pts[i+1][0]-pts[i][0]) + abs(pts[i+1][1]-pts[i][1])
            r = min(radius, seg_before // 2, seg_after // 2)
            bx, by = pts[i][0], pts[i][1]
            if prev_dir == 'R': bx -= r
            elif prev_dir == 'L': bx += r
            elif prev_dir == 'D': by -= r
            elif prev_dir == 'U': by += r
            path_cmds.append(f'<a:lnTo><a:pt x="{bx}" y="{by}"/></a:lnTo>')
            st, sw = ARC[(prev_dir, next_dir)]
            path_cmds.append(f'<a:arcTo wR="{r}" hR="{r}" stAng="{st}" swAng="{sw}"/>')
        else:
            path_cmds.append(f'<a:lnTo><a:pt x="{pts[i][0]}" y="{pts[i][1]}"/></a:lnTo>')

    path_data = ''.join(path_cmds)
    head_xml = f'<a:headEnd type="{head}" w="med" len="med"/>' if head != "none" else '<a:headEnd type="none"/>'
    tail_xml = f'<a:tailEnd type="{tail}" w="med" len="med"/>' if tail != "none" else '<a:tailEnd type="none"/>'

    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{bbox_w}" cy="{bbox_h}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:ahLst/>
      <a:cxnLst/>
      <a:rect l="0" t="0" r="{bbox_w}" b="{bbox_h}"/>
      <a:pathLst>
        <a:path w="{bbox_w}" h="{bbox_h}">
          {path_data}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="{dash}"/>
      <a:round/>
      {head_xml}
      {tail_xml}
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # Background
    bg_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="F8FAFF"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="EEF2FA"/></a:gs>
          </a:gsLst>
          <a:lin ang="5400000" scaled="0"/>
        </a:gradFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
    </p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml.strip()))
    sid += 1

    # Title
    title_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Title"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="457200" y="228600"/><a:ext cx="8229600" cy="457200"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="2400" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="1F3864"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{TITLE}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''
    spTree.append(etree.fromstring(title_xml.strip()))
    sid += 1

    # Compute shape positions
    shape_rects = []
    for i, (text, preset, fill, border) in enumerate(STEPS):
        x = LEFT + i * (SW + GAP)
        cy = DIMS[preset]
        y = BASE_Y + (BASE_CY - cy) // 2
        shape_rects.append((x, y, SW, cy))

    # Arrows (render before shapes for z-order)
    for i in range(N - 1):
        ax = shape_rects[i][0] + SW
        ay = BASE_Y + BASE_CY // 2
        bx = shape_rects[i + 1][0]
        xml = build_routed_connector_xml(
            [(ax, ay), (bx, ay)], sid, f"Arrow {sid}",
            color="4472C4", width=25400, tail="triangle")
        spTree.append(etree.fromstring(xml))
        sid += 1

    # Shapes (on top of arrows)
    for i, (text, preset, fill, border) in enumerate(STEPS):
        x, y, cx, cy = shape_rects[i]
        xml = make_shape_xml(sid, text, text, preset, x, y, cx, cy, fill, border)
        spTree.append(etree.fromstring(xml.strip()))
        sid += 1

    # "YES" label near Review decision
    review_x = shape_rects[2][0]
    label_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="YesLabel"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{review_x + SW + GAP // 4}" y="{BASE_Y - 228600}"/><a:ext cx="457200" cy="342900"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="1200" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="27AE60"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>YES</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''
    spTree.append(etree.fromstring(label_xml.strip()))

    output = "flowchart.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")


if __name__ == "__main__":
    main()
