#!/usr/bin/env python3
"""
Org Chart — 3-level hierarchy (1 CEO, 3 VPs, 6 Leads) with gradient cards.
Run: pip install python-pptx lxml && python org_chart.py && open org_chart.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Hierarchy — reporting structure, top-down authority.
# Layout algorithm: Layered/hierarchical, top-to-bottom. Each level occupies a
#   horizontal row; children are centered beneath their parent.
# Gestalt similarity: Level-based color gradient — CEO (dark navy), VPs (blue),
#   Leads (green). Color encodes rank without the viewer needing to trace edges.
# Gestalt proximity: Intra-level gap (LEVEL_GAP) < inter-level gap (V_GAP),
#   reinforcing that same-level peers are a group.
# Scan path: Top-to-bottom, center-outward. CEO is the entry point (top-center,
#   darkest color, largest visual weight).
# Progressive disclosure: For deeper orgs (>3 levels or >6 nodes per level),
#   split into an overview slide (levels 0-1) and detail slides per VP subtree.
# Data-ink: Luminance modulation gradient (lumMod 75%) adds depth without a second
#   hex color. Shadows are subtle (30% alpha) for figure-ground only.

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import os

SLIDE_W, SLIDE_H = 9144000, 5143500

# --- Data ---
TITLE = "Figure 2: Organization Chart"
PEOPLE = [
    # (id, name, title, level, parent)
    ("ceo", "Jane Smith",    "Chief Executive Officer", 0, None),
    ("vp1", "Alex Johnson",  "VP Engineering",          1, "ceo"),
    ("vp2", "Maria Garcia",  "VP Marketing",            1, "ceo"),
    ("vp3", "David Chen",    "VP Operations",           1, "ceo"),
    ("l1",  "Sam Lee",       "Frontend Lead",           2, "vp1"),
    ("l2",  "Priya Patel",   "Backend Lead",            2, "vp1"),
    ("l3",  "Tom Brown",     "Content Lead",            2, "vp2"),
    ("l4",  "Ana Silva",     "Growth Lead",             2, "vp2"),
    ("l5",  "Kevin Park",    "DevOps Lead",             2, "vp3"),
    ("l6",  "Lisa Wong",     "Support Lead",            2, "vp3"),
]

LEVEL_COLORS = {
    0: ("1F3864", "16294A"),  # dark navy
    1: ("4472C4", "2E4FA3"),  # blue
    2: ("70AD47", "507C32"),  # green
}

# --- Layout ---
SH = 571500          # shape height
V_GAP = 685800       # vertical gap between levels
MARGIN_TOP = 685800
LEVEL_SW = {0: 1600200, 1: 1600200, 2: 1280160}
LEVEL_GAP = {0: 304800, 1: 304800, 2: 182880}


def center_row(n, sw, gap):
    total = n * sw + (n - 1) * gap
    left = (SLIDE_W - total) // 2
    return [left + i * (sw + gap) for i in range(n)]


def make_card_xml(sid, key, name, title, x, y, sw, fill, border):
    """Rounded rect card with gradient, shadow, two-line text."""
    return f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{sid}" name="{key}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{sw}" cy="{SH}"/></a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="{fill}"><a:lumMod val="75000"/></a:srgbClr></a:gs>
          </a:gsLst>
          <a:lin ang="5400000" scaled="0"/>
        </a:gradFill>
        <a:ln w="19050"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
        <a:effectLst>
          <a:outerShdw blurRad="38100" dist="25400" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="000000"><a:alpha val="30000"/></a:srgbClr>
          </a:outerShdw>
        </a:effectLst>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="1100" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{name}</a:t></a:r>
        </a:p>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="900" b="0" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"><a:alpha val="80000"/></a:srgbClr></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{title}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''


def build_routed_connector_xml(waypoints, sid, name="Connector",
                                color="888888", width=19050,
                                dash="solid", tail="triangle",
                                head="none", radius=150000):
    """Orthogonal connector with curved elbows using custom geometry."""
    xs = [p[0] for p in waypoints]
    ys = [p[1] for p in waypoints]
    min_x, min_y = min(xs), min(ys)
    max_x, max_y = max(xs), max(ys)
    bbox_w = max(max_x - min_x, 1)
    bbox_h = max(max_y - min_y, 1)

    pts = [(x - min_x, y - min_y) for x, y in waypoints]

    def direction(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        if dx > 0: return 'R'
        if dx < 0: return 'L'
        if dy > 0: return 'D'
        return 'U'

    ARC = {
        ('R', 'D'): (16200000, 5400000),   ('R', 'U'): (5400000, -5400000),
        ('L', 'D'): (16200000, -5400000),  ('L', 'U'): (5400000, 5400000),
        ('D', 'R'): (10800000, -5400000),  ('D', 'L'): (0, 5400000),
        ('U', 'R'): (10800000, 5400000),   ('U', 'L'): (0, -5400000),
    }

    path_cmds = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for i in range(1, len(pts)):
        if i < len(pts) - 1:
            prev_dir = direction(pts[i-1], pts[i])
            next_dir = direction(pts[i], pts[i+1])
            seg_before = abs(pts[i][0]-pts[i-1][0]) + abs(pts[i][1]-pts[i-1][1])
            seg_after = abs(pts[i+1][0]-pts[i][0]) + abs(pts[i+1][1]-pts[i][1])
            r = min(radius, seg_before // 2, seg_after // 2)
            bx, by = pts[i][0], pts[i][1]
            if prev_dir == 'R': bx -= r
            elif prev_dir == 'L': bx += r
            elif prev_dir == 'D': by -= r
            elif prev_dir == 'U': by += r
            path_cmds.append(f'<a:lnTo><a:pt x="{bx}" y="{by}"/></a:lnTo>')
            st, sw = ARC[(prev_dir, next_dir)]
            path_cmds.append(f'<a:arcTo wR="{r}" hR="{r}" stAng="{st}" swAng="{sw}"/>')
        else:
            path_cmds.append(f'<a:lnTo><a:pt x="{pts[i][0]}" y="{pts[i][1]}"/></a:lnTo>')

    path_data = ''.join(path_cmds)
    head_xml = f'<a:headEnd type="{head}" w="med" len="med"/>' if head != "none" else '<a:headEnd type="none"/>'
    tail_xml = f'<a:tailEnd type="{tail}" w="med" len="med"/>' if tail != "none" else '<a:tailEnd type="none"/>'

    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{bbox_w}" cy="{bbox_h}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:ahLst/>
      <a:cxnLst/>
      <a:rect l="0" t="0" r="{bbox_w}" b="{bbox_h}"/>
      <a:pathLst>
        <a:path w="{bbox_w}" h="{bbox_h}">
          {path_data}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="{dash}"/>
      <a:round/>
      {head_xml}
      {tail_xml}
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # Background gradient
    bg_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="F8FAFF"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="EEF2FA"/></a:gs>
          </a:gsLst>
          <a:lin ang="5400000" scaled="0"/>
        </a:gradFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
    </p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml.strip()))
    sid += 1

    # Title
    title_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Title"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="457200" y="152400"/><a:ext cx="8229600" cy="400050"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="2200" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="1F3864"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{TITLE}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''
    spTree.append(etree.fromstring(title_xml.strip()))
    sid += 1

    # Compute positions by level
    levels = {}
    for pid, name, title, level, parent in PEOPLE:
        levels.setdefault(level, []).append((pid, name, title, parent))

    positions = {}  # pid -> (x, y, sw)
    for lvl in sorted(levels.keys()):
        people = levels[lvl]
        sw = LEVEL_SW[lvl]
        gap = LEVEL_GAP[lvl]
        xs = center_row(len(people), sw, gap)
        y = MARGIN_TOP + lvl * (SH + V_GAP)
        for i, (pid, name, title, parent) in enumerate(people):
            positions[pid] = (xs[i], y, sw)

    # Lines (render before cards for z-order) — orthogonal routing
    for pid, name, title, level, parent in PEOPLE:
        if parent is None:
            continue
        px, py, psw = positions[parent]
        cx, cy, csw = positions[pid]
        x1 = px + psw // 2   # parent bottom center
        y1 = py + SH
        x2 = cx + csw // 2   # child top center
        y2 = cy
        mid_y = y1 + (y2 - y1) // 2
        if x1 == x2:
            # Vertically aligned — straight down
            waypoints = [(x1, y1), (x2, y2)]
        else:
            # Z-shaped: down, across, down
            waypoints = [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]
        xml = build_routed_connector_xml(
            waypoints, sid, f"Line {sid}",
            color="AAAAAA", width=19050, tail="triangle", radius=120000)
        spTree.append(etree.fromstring(xml))
        sid += 1

    # Cards (on top)
    for pid, name, title, level, parent in PEOPLE:
        x, y, sw = positions[pid]
        fill, border = LEVEL_COLORS[level]
        xml = make_card_xml(sid, pid, name, title, x, y, sw, fill, border)
        spTree.append(etree.fromstring(xml.strip()))
        sid += 1

    output = "org_chart.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")


if __name__ == "__main__":
    main()
