#!/usr/bin/env python3
"""
M365 GenAI Chat System Architecture — Hub-and-spoke layout.
Run: pip install python-pptx lxml && python m365_genai_architecture.py && open m365_genai_architecture.pptx
"""

# DESIGN RATIONALE:
# Visual rhetoric: Hub-and-spoke — Copilot Studio Agent at center, 5 functional
#   regions surrounding it (Delivery, Prompt Library, Doc Mgmt, Retrieval, Governance).
# Layout algorithm: Central hub with surrounding regions in specific grid positions.
# Scan path: Top (delivery) → center (agent) → sides (prompt lib / doc mgmt) → bottom (retrieval, governance).
# Gestalt common region: Region background rects group functional areas.
# Gestalt similarity: Color encodes functional area (blue, orange, green, purple, gray).
# Data-ink: Solid lines = primary flow, dashed = secondary/supporting.

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
import os

# ---- Slide Dimensions (widescreen 13.33" x 7.5") ----
SLIDE_W = 12192000  # 13.33"
SLIDE_H = 6858000   # 7.5"

TITLE = "M365 GenAI Chat System Architecture"

# ---- Helper: inches to EMU ----
def IN(inches):
    return int(inches * 914400)

# ---- Style Colors: (fill, border) ----
STYLE_COLORS = {
    "interface":     ("4472C4", "2E4FA3"),
    "interface_hub": ("2E86C1", "1A5276"),
    "search":        ("ED7D31", "C05C13"),
    "docmgmt":       ("27AE60", "1A7A41"),
    "docmgmt_store": ("1A7A41", "0E5028"),
    "prompt":        ("7B68EE", "5B48CE"),
    "governance":    ("5D6D7E", "2C3E50"),
}

# ---- Region Definitions: (label, x, y, w, h, fill, border) ----
REGIONS = [
    ("DELIVERY &amp; INTERFACE",  IN(0.15), IN(0.40), IN(13.03), IN(0.90), "E8EAF0", "C5CBD8"),
    ("SHARED PROMPT LIBRARY", IN(0.15), IN(1.45), IN(2.80),  IN(3.20), "E8F5E9", "A5D6A7"),
    ("DOCUMENT MANAGEMENT",   IN(7.68), IN(1.45), IN(5.50),  IN(3.20), "E8F5E9", "A5D6A7"),
    ("RETRIEVAL &amp; SEARCH",    IN(3.10), IN(4.80), IN(6.00),  IN(1.30), "F1F8E9", "C5E1A5"),
    ("GOVERNANCE",            IN(0.15), IN(6.25), IN(13.03), IN(0.90), "E8EAF0", "C5CBD8"),
]

# ---- Component Definitions ----
# (id, line1, line2, style, x_in, y_in, w_in, h_in)
COMPONENTS = [
    # Delivery & Interface
    ("teams",       "Microsoft Teams",           "(1:1 Chat)",                "interface",     2.50, 0.53, 1.70, 0.65),
    ("powerapps",   "Power Apps Web Client",     "(GCC High Compatibility)",  "interface",     9.10, 0.53, 1.70, 0.65),

    # Central Hub (floating)
    ("agent",       "Copilot Studio Agent",      "(Per-Team/Project)",        "interface_hub", 5.00, 2.10, 2.30, 0.90),
    ("gen_answers", "Generative Answers",         "Node",                     "docmgmt",       5.30, 3.25, 1.70, 0.60),

    # Shared Prompt Library
    ("lists",       "Microsoft Lists",           "(Backend - PnP Gallery View)", "prompt",     0.40, 1.85, 2.30, 0.65),
    ("approval",    "Power Automate",            "Approval Flow (Curation)",     "prompt",     0.40, 3.55, 2.30, 0.65),

    # Document Management
    ("sp_docs",     "SharePoint Online",         "Document Libraries",        "docmgmt",       7.90, 1.85, 2.10, 0.70),
    ("pauto_cloud", "Power Automate",            "Cloud Flow",                "docmgmt",      10.30, 1.85, 1.70, 0.70),
    ("ai_builder",  "AI Builder",                "(OCR, Doc Processing)",     "docmgmt",       7.90, 3.00, 1.70, 0.70),
    ("sp_list",     "SharePoint List /",         "Dataverse Table",           "docmgmt_store", 9.80, 3.00, 1.40, 0.70),
    ("dv_metadata", "Dataverse Tables",          "(Metadata &amp; Audit)",        "docmgmt_store",11.40, 3.00, 1.50, 0.70),

    # Retrieval & Search
    ("graph_ground","Tenant Graph Grounding",    "(M365 Semantic Index)",     "search",        3.30, 4.95, 2.50, 0.65),
    ("sp_search",   "SharePoint Search Index",   "(BM25 Keyword Search)",     "search",        6.20, 4.95, 2.50, 0.65),

    # Governance (horizontal chain)
    ("entra",       "Microsoft",                 "Entra ID",                  "governance",    1.20, 6.43, 1.30, 0.55),
    ("gov_sp",      "SharePoint",                "",                          "governance",    3.50, 6.43, 1.30, 0.55),
    ("gov_dv",      "Dataverse",                 "(Audit Logs)",              "governance",    5.80, 6.43, 1.30, 0.55),
    ("teams_admin", "Teams Admin",               "Center",                    "governance",    8.10, 6.43, 1.30, 0.55),
    ("gov_teams",   "Microsoft",                 "Teams",                     "governance",   10.40, 6.43, 1.30, 0.55),
]

# Build position lookup: id -> (x_emu, y_emu, w_emu, h_emu)
POS = {}
COMP_LOOKUP = {}
for comp in COMPONENTS:
    cid, l1, l2, style, xi, yi, wi, hi = comp
    POS[cid] = (IN(xi), IN(yi), IN(wi), IN(hi))
    COMP_LOOKUP[cid] = comp

# ---- Arrows: (src, tgt, label, color, dash) ----
ARROWS = [
    ("teams",       "agent",       "Side-load /\nAdmin Push",       "4472C4", "solid"),
    ("powerapps",   "agent",       "Fallback\nInterface",           "4472C4", "dash"),
    ("agent",       "sp_docs",     "Knowledge\nSource",             "27AE60", "solid"),
    ("agent",       "gen_answers", "",                              "2E86C1", "solid"),
    ("agent",       "graph_ground","Copilot Studio\nAgent-Team",    "ED7D31", "solid"),
    ("lists",       "agent",       "Agent-Level\nKnowledge Source",  "7B68EE", "dash"),
    ("lists",       "approval",    "Prompt\nSubmission",            "7B68EE", "solid"),
    ("sp_docs",     "pauto_cloud", "Upload\nTrigger",               "27AE60", "solid"),
    ("pauto_cloud", "ai_builder",  "",                              "27AE60", "solid"),
    ("ai_builder",  "sp_list",     "Extracted\nFields",             "27AE60", "solid"),
    ("sp_list",     "dv_metadata", "Structured\nResults",           "27AE60", "solid"),
    ("sp_docs",     "sp_search",   "SharePoint\nSearch Index",      "ED7D31", "dash"),
    ("entra",       "gov_sp",      "Delegated\nPermissions",        "5D6D7E", "solid"),
    ("gov_sp",      "gov_dv",      "Replication",                   "5D6D7E", "solid"),
    ("gov_dv",      "teams_admin", "Interactions &amp;\nProcessing",    "5D6D7E", "solid"),
    ("teams_admin", "gov_teams",   "App\nDeployment",               "5D6D7E", "solid"),
]

# ---- Icon Map: component_id -> Iconify key ----
# Governance row excluded (shapes too small for icons at 1.30" × 0.55")
ICONS = {
    "teams":        "mdi:microsoft-teams",
    "powerapps":    "fluent:apps-list-24-filled",
    "agent":        "fluent:bot-sparkle-24-filled",
    "gen_answers":  "fluent:sparkle-24-filled",
    "lists":        "fluent:apps-list-detail-24-filled",
    "approval":     "fluent:cloud-flow-24-filled",
    "sp_docs":      "mdi:microsoft-sharepoint",
    "pauto_cloud":  "fluent:cloud-flow-24-filled",
    "ai_builder":   "fluent:brain-circuit-24-filled",
    "sp_list":      "mdi:database",
    "dv_metadata":  "mdi:database",
    "graph_ground": "fluent:plug-connected-24-filled",
    "sp_search":    "mdi:magnify",
}

# Per-shape icon sizing overrides: component_id -> (icon_size_inches, tIns_emu, top_offset_emu)
ICON_SIZE_OVERRIDES = {
    "agent":       (0.35, 365760, IN(0.10)),   # Hub: 2.30" × 0.90"
    "gen_answers": (0.22, 246888, IN(0.04)),   # Compact: 1.70" × 0.60"
}
# Default for standard shapes (0.65-0.70" tall)
DEFAULT_ICON_SIZING = (0.25, 274320, IN(0.06))


# ---- Edge anchor helpers ----
def _rc(p):
    """Right center of shape."""
    return (p[0] + p[2], p[1] + p[3] // 2)

def _lc(p):
    """Left center of shape."""
    return (p[0], p[1] + p[3] // 2)

def _bc(p):
    """Bottom center of shape."""
    return (p[0] + p[2] // 2, p[1] + p[3])

def _tc(p):
    """Top center of shape."""
    return (p[0] + p[2] // 2, p[1])

def _ccx(p):
    """Center X of shape."""
    return p[0] + p[2] // 2

def _ccy(p):
    """Center Y of shape."""
    return p[1] + p[3] // 2


# ---- Routing ----
def compute_waypoints(src_id, tgt_id):
    """Compute orthogonal waypoints for a connector."""
    sp, tp = POS[src_id], POS[tgt_id]
    s_cx, s_cy = _ccx(sp), _ccy(sp)
    t_cx, t_cy = _ccx(tp), _ccy(tp)

    same_row = abs(s_cy - t_cy) < IN(0.4)

    # Delivery → Agent: Z-shape vertical down
    if src_id == "teams" and tgt_id == "agent":
        x1, y1 = _bc(sp)
        x2, y2 = _tc(tp)
        mid_y = (y1 + y2) // 2
        return [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]

    if src_id == "powerapps" and tgt_id == "agent":
        x1, y1 = _bc(sp)
        x2, y2 = _tc(tp)
        mid_y = (y1 + y2) // 2
        return [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]

    # Agent → gen_answers: short vertical down
    if src_id == "agent" and tgt_id == "gen_answers":
        return [_bc(sp), _tc(tp)]

    # Agent → graph_ground: Z-shape vertical down
    if src_id == "agent" and tgt_id == "graph_ground":
        x1, y1 = _bc(sp)
        x2, y2 = _tc(tp)
        mid_y = (y1 + y2) // 2
        return [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]

    # Agent → sp_docs: horizontal right with Z-route
    if src_id == "agent" and tgt_id == "sp_docs":
        sx, sy = _rc(sp)
        tx, ty = _lc(tp)
        mid_x = (sx + tx) // 2
        return [(sx, sy), (mid_x, sy), (mid_x, ty), (tx, ty)]

    # Lists → Agent: horizontal right (across gap)
    if src_id == "lists" and tgt_id == "agent":
        sx, sy = _rc(sp)
        tx, ty = _lc(tp)
        mid_x = (sx + tx) // 2
        return [(sx, sy), (mid_x, sy), (mid_x, ty), (tx, ty)]

    # Lists → approval: vertical down
    if src_id == "lists" and tgt_id == "approval":
        return [_bc(sp), _tc(tp)]

    # sp_docs → pauto_cloud: horizontal right
    if src_id == "sp_docs" and tgt_id == "pauto_cloud":
        return [_rc(sp), _lc(tp)]

    # pauto_cloud → ai_builder: Z-shape down
    if src_id == "pauto_cloud" and tgt_id == "ai_builder":
        x1, y1 = _bc(sp)
        x2, y2 = _tc(tp)
        mid_y = (y1 + y2) // 2
        return [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]

    # ai_builder → sp_list: horizontal right
    if src_id == "ai_builder" and tgt_id == "sp_list":
        return [_rc(sp), _lc(tp)]

    # sp_list → dv_metadata: horizontal right
    if src_id == "sp_list" and tgt_id == "dv_metadata":
        return [_rc(sp), _lc(tp)]

    # sp_docs → sp_search: Z-shape vertical down
    if src_id == "sp_docs" and tgt_id == "sp_search":
        x1, y1 = _bc(sp)
        x2, y2 = _tc(tp)
        mid_y = (y1 + y2) // 2
        return [(x1, y1), (x1, mid_y), (x2, mid_y), (x2, y2)]

    # Governance chain: horizontal
    if same_row and s_cx < t_cx:
        return [_rc(sp), _lc(tp)]
    if same_row and s_cx >= t_cx:
        return [_lc(sp), _rc(tp)]

    # Generic Z-shape fallback
    if s_cx < t_cx:
        x1, y1 = _rc(sp)
        x2, y2 = _lc(tp)
    else:
        x1, y1 = _lc(sp)
        x2, y2 = _rc(tp)
    mid_x = (x1 + x2) // 2
    return [(x1, y1), (mid_x, y1), (mid_x, y2), (x2, y2)]


# ---- XML Builders ----

def build_routed_connector_xml(waypoints, sid, name="Connector",
                                color="888888", width=19050,
                                dash="solid", tail="triangle",
                                head="none", radius=100000):
    """Orthogonal connector with curved elbows using custom geometry."""
    xs = [p[0] for p in waypoints]
    ys = [p[1] for p in waypoints]
    min_x, min_y = min(xs), min(ys)
    max_x, max_y = max(xs), max(ys)
    bbox_w = max(max_x - min_x, 1)
    bbox_h = max(max_y - min_y, 1)

    pts = [(x - min_x, y - min_y) for x, y in waypoints]

    def direction(a, b):
        dx, dy = b[0] - a[0], b[1] - a[1]
        if dx > 0: return 'R'
        if dx < 0: return 'L'
        if dy > 0: return 'D'
        return 'U'

    ARC = {
        ('R', 'D'): (16200000, 5400000),   ('R', 'U'): (5400000, -5400000),
        ('L', 'D'): (16200000, -5400000),  ('L', 'U'): (5400000, 5400000),
        ('D', 'R'): (10800000, -5400000),  ('D', 'L'): (0, 5400000),
        ('U', 'R'): (10800000, 5400000),   ('U', 'L'): (0, -5400000),
    }

    path_cmds = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for i in range(1, len(pts)):
        if i < len(pts) - 1:
            prev_dir = direction(pts[i-1], pts[i])
            next_dir = direction(pts[i], pts[i+1])
            seg_before = abs(pts[i][0]-pts[i-1][0]) + abs(pts[i][1]-pts[i-1][1])
            seg_after = abs(pts[i+1][0]-pts[i][0]) + abs(pts[i+1][1]-pts[i][1])
            r = min(radius, seg_before // 2, seg_after // 2)
            bx, by = pts[i][0], pts[i][1]
            if prev_dir == 'R': bx -= r
            elif prev_dir == 'L': bx += r
            elif prev_dir == 'D': by -= r
            elif prev_dir == 'U': by += r
            path_cmds.append(f'<a:lnTo><a:pt x="{bx}" y="{by}"/></a:lnTo>')
            st, sw = ARC[(prev_dir, next_dir)]
            path_cmds.append(f'<a:arcTo wR="{r}" hR="{r}" stAng="{st}" swAng="{sw}"/>')
        else:
            path_cmds.append(f'<a:lnTo><a:pt x="{pts[i][0]}" y="{pts[i][1]}"/></a:lnTo>')

    path_data = ''.join(path_cmds)
    head_xml = f'<a:headEnd type="{head}" w="med" len="med"/>' if head != "none" else '<a:headEnd type="none"/>'
    tail_xml = f'<a:tailEnd type="{tail}" w="med" len="med"/>' if tail != "none" else '<a:tailEnd type="none"/>'

    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{min_x}" y="{min_y}"/><a:ext cx="{bbox_w}" cy="{bbox_h}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/>
      <a:gdLst/>
      <a:ahLst/>
      <a:cxnLst/>
      <a:rect l="0" t="0" r="{bbox_w}" b="{bbox_h}"/>
      <a:pathLst>
        <a:path w="{bbox_w}" h="{bbox_h}">
          {path_data}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{width}">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:prstDash val="{dash}"/>
      <a:round/>
      {head_xml}
      {tail_xml}
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
</p:sp>'''


def make_component_xml(sid, name, line1, line2, x, y, cx, cy, fill, border, font_size=900):
    """Create a rounded rectangle component with 1-2 line label."""
    text_paras = f'''<a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="{font_size}" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{line1}</a:t></a:r>
        </a:p>'''
    if line2:
        sub_size = max(font_size - 100, 700)
        text_paras += f'''
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="{sub_size}" b="0" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{line2}</a:t></a:r>
        </a:p>'''

    return f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr>
        <p:cNvPr id="{sid}" name="{name}"/>
        <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="{fill}"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="{border}"/></a:gs>
          </a:gsLst>
          <a:lin ang="5400000" scaled="0"/>
        </a:gradFill>
        <a:ln w="19050"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill><a:round/></a:ln>
        <a:effectLst>
          <a:outerShdw blurRad="40000" dist="25400" dir="5400000" algn="tl" rotWithShape="0">
            <a:srgbClr val="000000"><a:alpha val="30000"/></a:srgbClr>
          </a:outerShdw>
        </a:effectLst>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr" lIns="45720" rIns="45720" tIns="27432" bIns="27432">
          <a:normAutofit/>
        </a:bodyPr>
        <a:lstStyle/>
        {text_paras}
      </p:txBody>
    </p:sp>'''


def make_label_xml(sid, text, x, y, cx=None, cy=None, color="666666",
                   font_sz=700, bold=True, italic=False):
    """Transparent text box for connector labels or annotations."""
    if cx is None:
        cx = IN(1.20)
    if cy is None:
        cy = IN(0.55)
    b_attr = "1" if bold else "0"
    i_attr = "1" if italic else "0"
    # Support multiline text
    lines = text.split('\n')
    paras = ''
    for line in lines:
        paras += f'''<a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="{font_sz}" b="{b_attr}" i="{i_attr}" dirty="0">
            <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{line}</a:t></a:r>
        </a:p>'''

    return f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Lbl {sid}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:noFill/><a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        {paras}
      </p:txBody>
    </p:sp>'''


def compute_label_pos(waypoints, src_id, tgt_id):
    """Compute label position offset from the connector midpoint, avoiding shape overlap."""
    if len(waypoints) < 2:
        return waypoints[0][0], waypoints[0][1]

    # Use the midpoint segment
    mid_idx = len(waypoints) // 2
    ax, ay = waypoints[mid_idx - 1]
    bx, by = waypoints[mid_idx]
    mx = (ax + bx) // 2
    my = (ay + by) // 2

    # Determine if the mid-segment is horizontal or vertical
    is_horizontal = abs(bx - ax) > abs(by - ay)

    if is_horizontal:
        # Place label above the horizontal segment
        return mx - IN(0.55), my - IN(0.55)
    else:
        # Place label to the right of the vertical segment
        return mx + IN(0.10), my - IN(0.25)


# ---- Icon Helpers ----

def resolve_icon(icon_key, size=64, tint="FFFFFF", cache_dir="/tmp/diagram_icons"):
    """Resolve an icon key to a local PNG path via Iconify API."""
    os.makedirs(cache_dir, exist_ok=True)

    if "/" in icon_key or "\\" in icon_key or icon_key.endswith((".png", ".svg")):
        if icon_key.endswith(".png"):
            return icon_key if os.path.exists(icon_key) else None
        if icon_key.endswith(".svg"):
            if not os.path.exists(icon_key):
                return None
            try:
                import cairosvg
                png_path = os.path.join(cache_dir, os.path.basename(icon_key) + ".png")
                with open(icon_key, "r") as f:
                    svg_data = f.read()
                if tint:
                    svg_data = svg_data.replace("currentColor", f"#{tint}")
                cairosvg.svg2png(bytestring=svg_data.encode(), write_to=png_path,
                                 output_width=size, output_height=size)
                return png_path
            except ImportError:
                print(f"  Warning: cairosvg not installed, cannot convert {icon_key}")
                return None
        return None

    if ":" in icon_key:
        tint_suffix = f"_{tint}" if tint else "_orig"
        cache_name = f"{icon_key.replace(':', '_')}_{size}{tint_suffix}.png"
        png_path = os.path.join(cache_dir, cache_name)
        if os.path.exists(png_path):
            return png_path
        try:
            import cairosvg
            import urllib.request
            prefix, name = icon_key.split(":", 1)
            url = f"https://api.iconify.design/{prefix}/{name}.svg?width={size}&height={size}"
            req = urllib.request.Request(url, headers={"User-Agent": "python-pptx-diagrams/1.0"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                svg_data = resp.read().decode("utf-8")
            if tint:
                svg_data = svg_data.replace("currentColor", f"#{tint}")
            cairosvg.svg2png(bytestring=svg_data.encode(), write_to=png_path,
                             output_width=size, output_height=size)
            return png_path
        except ImportError:
            print(f"  Warning: cairosvg not installed, skipping icon '{icon_key}'")
            return None
        except Exception as e:
            print(f"  Warning: Failed to resolve icon '{icon_key}': {e}")
            return None

    return None


def enrich_shape_with_icon(slide, shape_name, icon_path,
                            icon_size_inches=0.25, tIns=274320, top_offset=None):
    """Add an icon image above the text in a named shape."""
    from pptx.util import Inches

    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break
    if target is None:
        print(f"  Warning: Shape '{shape_name}' not found")
        return False

    icon_w = icon_h = Inches(icon_size_inches)
    icon_left = target.left + (target.width - icon_w) // 2
    if top_offset is None:
        top_offset = Emu(91440)  # 0.1"
    icon_top = target.top + Emu(top_offset) if isinstance(top_offset, int) else target.top + top_offset

    try:
        slide.shapes.add_picture(icon_path, icon_left, icon_top, icon_w, icon_h)
    except Exception as e:
        print(f"  Warning: Failed to add icon to '{shape_name}': {e}")
        return False

    if target.has_text_frame:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        body_pr = target.text_frame._txBody.find(f"{{{ns}}}bodyPr")
        if body_pr is not None:
            body_pr.set("tIns", str(tIns))
            body_pr.set("anchor", "b")

    return True


def enrich_pptx_with_icons(pptx_path, icon_map, output_path=None,
                            icon_size=64, tint="FFFFFF", slide_index=0):
    """Post-process PPTX to add icons to named shapes with per-shape sizing."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    results = []

    for shape_name, icon_key in icon_map.items():
        icon_path = resolve_icon(icon_key, size=icon_size, tint=tint)
        if icon_path:
            sz, tIns, top_off = ICON_SIZE_OVERRIDES.get(shape_name, DEFAULT_ICON_SIZING)
            ok = enrich_shape_with_icon(slide, shape_name, icon_path,
                                         icon_size_inches=sz, tIns=tIns, top_offset=top_off)
            results.append((shape_name, ok))
        else:
            results.append((shape_name, False))

    prs.save(output_path or pptx_path)
    return results


# ---- Main ----

def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.shapes._spTree
    sid = 2

    # === 1. Background ===
    bg_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="Background"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="F5F6FA"/></a:solidFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>
    </p:sp>'''
    spTree.insert(2, etree.fromstring(bg_xml.strip()))
    sid += 1

    # === 2. Region backgrounds ===
    for label, rx, ry, rw, rh, fill, border in REGIONS:
        rgn_xml = f'''
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr><p:cNvPr id="{sid}" name="Region {label}"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
          <p:spPr>
            <a:xfrm><a:off x="{rx}" y="{ry}"/><a:ext cx="{rw}" cy="{rh}"/></a:xfrm>
            <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 5000"/></a:avLst></a:prstGeom>
            <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
            <a:ln w="12700"><a:solidFill><a:srgbClr val="{border}"/></a:solidFill></a:ln>
          </p:spPr>
          <p:txBody>
            <a:bodyPr wrap="square" anchor="t" lIns="91440" tIns="45720"><a:normAutofit/></a:bodyPr>
            <a:lstStyle/>
            <a:p><a:pPr algn="l"/>
              <a:r><a:rPr lang="en-US" sz="800" b="1" dirty="0">
                <a:solidFill><a:srgbClr val="888888"/></a:solidFill>
                <a:latin typeface="Calibri"/>
              </a:rPr><a:t>{label}</a:t></a:r>
            </a:p>
          </p:txBody>
        </p:sp>'''
        spTree.append(etree.fromstring(rgn_xml.strip()))
        sid += 1

    # === 3. Title bar ===
    title_xml = f'''
    <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:nvSpPr><p:cNvPr id="{sid}" name="TitleBar"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{IN(0.15)}" y="0"/><a:ext cx="{IN(13.03)}" cy="{IN(0.35)}"/></a:xfrm>
        <a:prstGeom prst="roundRect"><a:avLst><a:gd name="adj" fmla="val 8000"/></a:avLst></a:prstGeom>
        <a:gradFill>
          <a:gsLst>
            <a:gs pos="0"><a:srgbClr val="1F3864"/></a:gs>
            <a:gs pos="100000"><a:srgbClr val="16294A"/></a:gs>
          </a:gsLst>
          <a:lin ang="0" scaled="0"/>
        </a:gradFill>
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" anchor="ctr"><a:normAutofit/></a:bodyPr>
        <a:lstStyle/>
        <a:p><a:pPr algn="ctr"/>
          <a:r><a:rPr lang="en-US" sz="1400" b="1" dirty="0">
            <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
            <a:latin typeface="Calibri"/>
          </a:rPr><a:t>{TITLE}</a:t></a:r>
        </a:p>
      </p:txBody>
    </p:sp>'''
    spTree.append(etree.fromstring(title_xml.strip()))
    sid += 1

    # === 4. Connectors ===
    for src_id, tgt_id, label, color, dash in ARROWS:
        waypoints = compute_waypoints(src_id, tgt_id)
        xml = build_routed_connector_xml(
            waypoints, sid, f"Arrow {src_id}-{tgt_id}",
            color=color, width=19050, tail="stealth",
            dash=dash, radius=80000)
        spTree.append(etree.fromstring(xml))
        sid += 1

    # === 5. Connector labels (manually positioned to avoid overlaps) ===
    # (label_text, x_inches, y_inches, color, cx_inches, cy_inches)
    MANUAL_LABELS = [
        # teams → agent: on horizontal segment between delivery and hub
        ("Side-load /\nAdmin Push",      2.60, 1.35, "4472C4", 1.00, 0.55),
        # powerapps → agent: near powerapps, to the left
        ("Fallback\nInterface",          8.20, 1.35, "4472C4", 0.90, 0.55),
        # agent → sp_docs: above connector, left of sp_docs
        ("Knowledge\nSource",            7.40, 1.35, "27AE60", 0.80, 0.55),
        # agent → graph_ground: to the right of vertical drop
        ("Copilot Studio\nAgent-Team",   4.55, 3.90, "ED7D31", 1.10, 0.55),
        # lists → agent: above horizontal connector, clear of Side-load label
        ("Agent-Level\nKnowledge Source", 3.50, 1.80, "7B68EE", 1.20, 0.55),
        # lists → approval: to the right of vertical connector
        ("Prompt\nSubmission",           2.80, 2.60, "7B68EE", 0.90, 0.55),
        # sp_docs → pauto_cloud: above horizontal connector
        ("Upload\nTrigger",             10.60, 1.25, "27AE60", 0.85, 0.55),
        # ai_builder → sp_list: below horizontal connector
        ("Extracted\nFields",            9.60, 3.75, "27AE60", 0.80, 0.40),
        # sp_list → dv_metadata: below horizontal connector
        ("Structured\nResults",         10.60, 3.70, "27AE60", 0.90, 0.55),
        # sp_docs → sp_search: to the right of vertical drop
        ("SharePoint\nSearch Index",     9.20, 4.15, "ED7D31", 1.10, 0.55),
        # entra → gov_sp: above governance connector
        ("Delegated\nPermissions",       2.20, 5.90, "5D6D7E", 1.00, 0.55),
        # gov_sp → gov_dv: above governance connector
        ("Replication",                  4.55, 5.95, "5D6D7E", 0.90, 0.55),
        # gov_dv → teams_admin: above governance connector
        ("Interactions &amp;\nProcessing",6.90, 5.90, "5D6D7E", 1.00, 0.55),
        # teams_admin → gov_teams: above governance connector
        ("App\nDeployment",              9.20, 5.90, "5D6D7E", 0.90, 0.55),
    ]

    for label_text, lx, ly, color, lcx, lcy in MANUAL_LABELS:
        xml = make_label_xml(sid, label_text, IN(lx), IN(ly),
                             cx=IN(lcx), cy=IN(lcy),
                             color=color, font_sz=650)
        spTree.append(etree.fromstring(xml.strip()))
        sid += 1

    # === 6. Component shapes (on top) ===
    for comp in COMPONENTS:
        cid, line1, line2, style, xi, yi, wi, hi = comp
        x, y, cx, cy = POS[cid]
        fill, border = STYLE_COLORS[style]
        # Hub component gets larger font
        fs = 1000 if style == "interface_hub" else 900
        xml = make_component_xml(sid, cid, line1, line2, x, y, cx, cy,
                                  fill, border, font_size=fs)
        spTree.append(etree.fromstring(xml.strip()))
        sid += 1

    # === 7. Annotations ===
    # Scope isolation note inside Retrieval region bottom
    annot1_xml = make_label_xml(
        sid, "Scope Isolation: Separate Agent per Team/Project Site",
        IN(4.20), IN(5.65), cx=IN(3.80), cy=IN(0.28),
        color="5D6D7E", font_sz=550, bold=False, italic=True)
    spTree.append(etree.fromstring(annot1_xml.strip()))
    sid += 1

    # M365 Copilot License note near Graph Grounding
    annot2_xml = make_label_xml(
        sid, "Requires M365 Copilot License",
        IN(3.30), IN(4.60), cx=IN(2.50), cy=IN(0.35),
        color="ED7D31", font_sz=600, bold=False, italic=True)
    spTree.append(etree.fromstring(annot2_xml.strip()))
    sid += 1

    output = "m365_genai_architecture.pptx"
    prs.save(output)
    print(f"Saved: {os.path.abspath(output)}")

    # === 8. Icon enrichment (post-save) ===
    print("Adding icons...")
    results = enrich_pptx_with_icons(output, ICONS)
    added = sum(1 for _, ok in results if ok)
    failed = [(name, ok) for name, ok in results if not ok]
    print(f"Icons: {added}/{len(ICONS)} added successfully")
    if failed:
        for name, _ in failed:
            print(f"  Skipped: {name}")


if __name__ == "__main__":
    main()
